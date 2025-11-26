#!/usr/bin/env python3
"""
Add speaker notes to Tracepoint_Security_Architecture.pptx in-place.
Run from repo root:  python3 add_speaker_notes.py
"""

from pptx import Presentation

SOURCE = "Tracepoint_Security_Architecture.pptx"
TARGET = "Tracepoint_Security_Architecture_with_notes.pptx"

notes_by_index = {
    # Slide 0 – Cover
    0: (
        "Opening:\n"
        "- Tracepoint is our in-house platform built specifically for Swanson.\n"
        "- This deck explains the security architecture backing the CAPEX request.\n"
        "Key message: We are using the same security patterns as banks and large SaaS vendors."
    ),
    # Slide 1 – Authentication
    1: (
        "Authentication system:\n"
        "Bcrypt: Slow by design so brute-force guessing is impractical; even with leaked hashes,\n"
        "attackers need enormous compute to try passwords.\n"
        "10 salt rounds: The hash is recomputed 2^10 times; more rounds means more work for an attacker\n"
        "with almost no impact on a normal login.\n"
        "OAuth: We outsource identity checks to Google/Microsoft, so we never see those passwords.\n"
        "2FA/TOTP: Even if a password is stolen, an attacker still needs the one-time code from the\n"
        "user's device."
    ),
    # Slide 2 – Tokens
    2: (
        "Token-based sessions:\n"
        "Login token: Very short-lived, only used during the sign-in flow to minimize exposure.\n"
        "Access token: Carries permissions for API calls; it expires in about an hour to cap risk.\n"
        "Refresh token: Longer-lived, but can be revoked instantly if we see suspicious behavior.\n"
        "JWT: Digitally signed so any tampering is detected immediately and the token is rejected."
    ),
    # Slide 3 – Encryption
    3: (
        "Data encryption:\n"
        "AES-256 at rest: Same standard used by financial institutions; not realistically brute-forceable.\n"
        "IV per record: Even identical values encrypt differently, which hides patterns in the data.\n"
        "TLS 1.3 in transit: Protects against snooping and man-in-the-middle attacks on the network."
    ),
    # Slide 4 – RBAC
    4: (
        "Role-based access control (RBAC):\n"
        "Concept: We give people roles (Admin, Manager, Sales, etc.) and attach permissions to roles,\n"
        "not individuals. That makes access easier to reason about and audit.\n"
        "Granularity: 22 permission categories let us separate duties; for example, someone can export\n"
        "reports without being able to delete data or touch billing.\n"
        "Field-level control: We can allow a role to see a field like 'Revenue' but block edits, or\n"
        "hide sensitive fields such as salary from everyone except Finance."
    ),
    # Slide 5 – API Security
    5: (
        "API security and keys:\n"
        "Service accounts: Integrations use API keys instead of user passwords, so people can change jobs\n"
        "without breaking systems.\n"
        "Scoped keys: Each key only has the minimum permissions it needs, based on the attached role.\n"
        "Revocation: If a key leaks, we disable it centrally and all access stops immediately.\n"
        "Rate limiting: Protects our backend from abuse and helps absorb accidental or malicious spikes."
    ),
    # Slide 6 – Workspace Isolation
    6: (
        "Workspace isolation:\n"
        "Design: Each company's data lives in its own logical workspace, enforced at the database level.\n"
        "Effect: Even if an engineer or query is misconfigured, cross-workspace reads are blocked by design.\n"
        "Benefit: This is a strong guardrail against both coding mistakes and privilege misuse."
    ),
    # Slide 7 – Audit Trail & Monitoring
    7: (
        "Audit trail and monitoring:\n"
        "Coverage: We log logins, permission changes, exports, API key usage, and 2FA changes.\n"
        "Why it matters: If something looks odd - a spike in exports, a new admin role - we can answer\n"
        "who did what, from where, and when within seconds.\n"
        "Compliance: These logs support internal audits and external compliance requests."
    ),
    # Slide 8 – Attack Prevention
    8: (
        "Attack prevention:\n"
        "Rate limits: Hard caps on login attempts, password resets, and API calls make brute-force and\n"
        "denial-of-service style abuse much harder.\n"
        "CAPTCHA options: We can turn on CAPTCHA for high-risk flows like login and password reset to\n"
        "slow down automated attacks.\n"
        "File controls: Upload size limits and validation protect storage and reduce the risk from\n"
        "malicious files."
    ),
    # Slide 9 – Compliance & Standards
    9: (
        "Compliance and standards alignment:\n"
        "OWASP Top 10: The design explicitly addresses the top web risks, like injection and broken access\n"
        "control, using parameterized queries, strong auth, and strict permission checks.\n"
        "GDPR / SOC 2 / ISO 27001: We follow the patterns these frameworks expect: data minimization,\n"
        "access logging, and controlled change processes.\n"
        "Business translation: We are building Tracepoint in a way that will not become a compliance\n"
        "liability as the company grows."
    ),
    # Slide 10 – Security Posture
    10: (
        "Current security posture:\n"
        "Implemented: Password hashing, SSO, encryption, RBAC, isolation, audit logging, and rate limits\n"
        "are already in place; this is not a slide of future promises.\n"
        "Maturity: We are operating at a level appropriate for a modern SaaS used by industrial customers,\n"
        "with clear room to tighten further over time.\n"
        "Message to execs: This CAPEX is about continuing to invest and stay ahead of risks, not fixing\n"
        "a broken foundation."
    ),
    # Slide 11 – Summary / Closing
    11: (
        "Summary and closing:\n"
        "Core message: Tracepoint gives Swanson a platform where security is engineered in, not bolted on.\n"
        "Risk framing: The cost of building this correctly now is small compared to the potential impact\n"
        "of a breach or compliance failure later.\n"
        "Final line to emphasize: Swanson's data is protected by the same security standards used by banks\n"
        "and governments."
    ),
}


def main() -> None:
    prs = Presentation(SOURCE)

    for idx, text in notes_by_index.items():
        if idx >= len(prs.slides):
            continue
        slide = prs.slides[idx]
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = ""  # clear existing
        p = tf.paragraphs[0]
        p.text = text

    prs.save(TARGET)
    print(f"✅ Speaker notes added to {len(notes_by_index)} slides.")
    print(f"📁 Saved as: {TARGET}")


if __name__ == "__main__":
    main()
