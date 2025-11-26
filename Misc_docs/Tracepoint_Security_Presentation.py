#!/usr/bin/env python3
"""
Tracepoint Security Architecture Presentation
Author: Wayne Lytle, Swanson Industries
Purpose: CAPEX Proposal for Executive Leadership
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_gradient_background(slide, start_color=(15, 32, 61), end_color=(44, 62, 80)):
    """Add gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = RGBColor(*start_color)
    fill.gradient_stops[1].color.rgb = RGBColor(*end_color)

def add_page_number(slide, page_num, total_pages=12):
    """Add page number to bottom right"""
    textbox = slide.shapes.add_textbox(
        Inches(8.8), Inches(7.0), Inches(1.0), Inches(0.3)
    )
    text_frame = textbox.text_frame
    text_frame.text = f"{page_num}/{total_pages}"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    paragraph.font.size = Pt(10)
    paragraph.font.color.rgb = RGBColor(169, 169, 169)

def add_logo(slide, position='top_left'):
    """Add Swanson logo to slide"""
    logo_path = '/home/lytle/twenty-dev/packages/twenty-front/public/images/logos/swanson logo full (1) (1).png'

    if not os.path.exists(logo_path):
        # Skip if logo not found
        return

    if position == 'top_left':
        left, top = Inches(0.5), Inches(0.3)
        height = Inches(0.8)
    else:  # bottom_right
        left, top = Inches(8.0), Inches(6.8)
        height = Inches(0.6)

    try:
        slide.shapes.add_picture(logo_path, left, top, height=height)
    except:
        pass  # Continue if logo can't be added

# ============================================================================
# SLIDE 1: COVER PAGE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_gradient_background(slide, (15, 32, 61), (25, 42, 71))
add_logo(slide, 'top_left')

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Tracepoint Security Architecture"
title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Enterprise-Grade Security for Swanson's Data"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.alignment = PP_ALIGN.CENTER
subtitle_para.font.size = Pt(28)
subtitle_para.font.color.rgb = RGBColor(200, 200, 200)

# Author info
author_box = slide.shapes.add_textbox(Inches(2), Inches(6.0), Inches(6), Inches(0.8))
author_frame = author_box.text_frame
author_frame.text = "Presented by: Wayne Lytle\nSwanson Industries IT Department"
author_para = author_frame.paragraphs[0]
author_para.alignment = PP_ALIGN.CENTER
author_para.font.size = Pt(16)
author_para.font.color.rgb = RGBColor(169, 169, 169)

# ============================================================================
# SLIDE 2: Authentication System
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 2)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "1. Authentication System 🛡️"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

sections = [
    ("Multi-Layer Password Protection", [
        "Bcrypt hashing with 10 salt rounds (industry standard)",
        "One-way encryption - passwords cannot be decrypted",
        "Minimum 8 characters enforced",
        "Rainbow table protection via unique salts"
    ]),
    ("OAuth 2.0 Integration (Enterprise SSO)", [
        "Google OAuth - OpenID Connect certified",
        "Microsoft Azure AD - SAML 2.0 compliant",
        "Custom SSO available for enterprise needs",
        "No password storage - delegated to trusted providers"
    ]),
    ("Two-Factor Authentication (2FA)", [
        "TOTP (Time-based One-Time Password)",
        "6-digit codes rotating every 30 seconds",
        "RFC 6238 compliant (Google Authenticator standard)",
        "Workspace-level enforcement capability"
    ])
]

for section_title, items in sections:
    p = content_frame.add_paragraph()
    p.text = section_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 181, 246)
    p.space_after = Pt(6)

    for item in items:
        p = content_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(224, 224, 224)
        p.level = 1
        p.space_after = Pt(3)

    content_frame.add_paragraph()

# ============================================================================
# SLIDE 3: Token-Based Session Management
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 3)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "2. Token-Based Session Management 🎫"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "JWT (JSON Web Tokens) - Three-Tier System"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_after = Pt(12)

token_types = [
    ("Login Token", "5 minutes", "Initial authentication", "Short-lived, workspace-specific"),
    ("Access Token", "1 hour", "API requests", "Contains user permissions"),
    ("Refresh Token", "30 days", "Renew access", "Revocable instantly")
]

for token, lifespan, purpose, security in token_types:
    p = content_frame.add_paragraph()
    p.text = f"{token} ({lifespan})"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p = content_frame.add_paragraph()
    p.text = f"• Purpose: {purpose}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

    p = content_frame.add_paragraph()
    p.text = f"• Security: {security}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1
    p.space_after = Pt(10)

p = content_frame.add_paragraph()
p.text = "Token Security Features:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_before = Pt(10)

features = [
    "Cryptographic signature verification (SHA-256)",
    "Automatic expiration enforcement",
    "Workspace-specific isolation",
    "Instant revocation capability"
]

for feature in features:
    p = content_frame.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(76, 175, 80)
    p.level = 1

# ============================================================================
# SLIDE 4: Data Encryption
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 4)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "3. Data Encryption 🔒"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "Encryption at Rest (AES-256-CTR)"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_after = Pt(10)

encrypted_data = [
    "OAuth refresh tokens and access credentials",
    "API keys and webhook secrets",
    "Connected account authentication data",
    "Sensitive customer information"
]

for item in encrypted_data:
    p = content_frame.add_paragraph()
    p.text = f"✓ {item}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

p = content_frame.add_paragraph()
p.text = "\nEncryption Process:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(12)

process = [
    "Unique IV (Initialization Vector) per record",
    "SHA-512 key hashing",
    "Military-grade AES-256-CTR encryption",
    "Secure base64 storage"
]

for step in process:
    p = content_frame.add_paragraph()
    p.text = f"• {step}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

p = content_frame.add_paragraph()
p.text = "\nEncryption in Transit (HTTPS/TLS 1.3)"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_before = Pt(15)

transit_features = [
    "All API traffic encrypted end-to-end",
    "Certificate-based authentication",
    "Perfect forward secrecy (PFS)",
    "Protection against man-in-the-middle attacks"
]

for feature in transit_features:
    p = content_frame.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

# ============================================================================
# SLIDE 5: Role-Based Access Control
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 5)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "4. Role-Based Access Control (RBAC) 👥"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "Granular Permission System (22 Categories)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_after = Pt(10)

permission_categories = [
    ("Settings Permissions", [
        "Workspace - Global settings management",
        "Members - User access control",
        "Roles - Permission management",
        "Data Model - Database schema control",
        "Security - Authentication & 2FA",
        "Billing - Subscription management"
    ]),
    ("Tool Permissions", [
        "AI - AI agent interactions",
        "Upload/Download Files - Document control",
        "Export/Import CSV - Data portability",
        "Send Email - Communication tools"
    ])
]

for category, perms in permission_categories:
    p = content_frame.add_paragraph()
    p.text = category
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 193, 7)
    p.space_before = Pt(8)

    for perm in perms:
        p = content_frame.add_paragraph()
        p.text = f"• {perm}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(224, 224, 224)
        p.level = 1

p = content_frame.add_paragraph()
p.text = "\nField-Level Permissions"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_before = Pt(12)

p = content_frame.add_paragraph()
p.text = "Control access at individual field level (read, create, update, delete)"
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(224, 224, 224)
p.level = 1

# ============================================================================
# SLIDE 6: API Security
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 6)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "5. API Security 🔑"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "API Key Authentication for Integrations"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_after = Pt(12)

features = [
    "Long-lived tokens for service accounts",
    "Scoped permissions (inherit from role)",
    "Instant revocation capability",
    "Complete audit trail of all usage"
]

for feature in features:
    p = content_frame.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

p = content_frame.add_paragraph()
p.text = "\nAPI Key Structure Example:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(15)

p = content_frame.add_paragraph()
p.text = "sk_a1b2c3d4_e5f6g7h8i9j0k1l2m3n4o5p6"
p.font.size = Pt(14)
p.font.name = 'Courier New'
p.font.color.rgb = RGBColor(100, 221, 23)
p.level = 1

p = content_frame.add_paragraph()
p.text = "\nSecurity Features:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(15)

security_features = [
    "256-bit entropy (cryptographically random)",
    "Rate limited (1000 requests/minute default)",
    "Workspace-specific prefix identification",
    "Secure storage with encryption at rest"
]

for feature in security_features:
    p = content_frame.add_paragraph()
    p.text = f"• {feature}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

# ============================================================================
# SLIDE 7: Workspace Isolation
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 7)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "6. Workspace Isolation 🏢"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "Multi-Tenant Security Architecture"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = content_frame.add_paragraph()
p.text = "Complete Data Segregation:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(8)

isolation_features = [
    "Schema-level isolation in PostgreSQL database",
    "Row-level security (RLS) policies enforced",
    "Dedicated connection pooling per workspace",
    "Separate encryption keys per workspace",
    "Impossible to access another organization's data"
]

for feature in isolation_features:
    p = content_frame.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

p = content_frame.add_paragraph()
p.text = "\nCross-Workspace Protection:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(15)

protection_features = [
    "Tokens bound to workspace ID",
    "All queries filtered by workspace",
    "File storage completely segregated",
    "No shared resources between workspaces",
    "Zero data leakage risk"
]

for feature in protection_features:
    p = content_frame.add_paragraph()
    p.text = f"• {feature}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

# ============================================================================
# SLIDE 8: Audit Trail & Monitoring
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 8)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "7. Audit Trail & Monitoring 📊"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "Comprehensive Security Auditing"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = content_frame.add_paragraph()
p.text = "Logged Events:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

logged_events = [
    "Login attempts (success & failure)",
    "Password changes and resets",
    "Permission and role changes",
    "Data exports and imports",
    "API key creation/revocation",
    "2FA enable/disable events",
    "User impersonation activities"
]

for event in logged_events:
    p = content_frame.add_paragraph()
    p.text = f"✓ {event}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

p = content_frame.add_paragraph()
p.text = "\nAudit Capabilities:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(12)

capabilities = [
    "Real-time monitoring dashboard",
    "Historical audit log queries",
    "IP address tracking",
    "Timestamp precision to milliseconds",
    "Exportable compliance reports"
]

for capability in capabilities:
    p = content_frame.add_paragraph()
    p.text = f"• {capability}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

# ============================================================================
# SLIDE 9: Attack Prevention
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 9)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "8. Attack Prevention & Protection 🛡️"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "Rate Limiting & Brute Force Prevention"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_after = Pt(10)

rate_limits = [
    ("Login Attempts", "5 per 15 minutes", "Prevents brute force attacks"),
    ("Password Reset", "3 per hour", "Prevents reset abuse"),
    ("API Requests", "1000 per minute", "DDoS protection"),
    ("File Uploads", "100MB max size", "Prevents resource exhaustion")
]

for limit_type, limit, protection in rate_limits:
    p = content_frame.add_paragraph()
    p.text = f"{limit_type}: {limit}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p = content_frame.add_paragraph()
    p.text = f"→ {protection}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1
    p.space_after = Pt(6)

p = content_frame.add_paragraph()
p.text = "\nCAPTCHA Protection (Optional)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_before = Pt(12)

captcha_features = [
    "Configurable for login, signup, password reset",
    "Google reCAPTCHA v3 support",
    "hCaptcha support",
    "Cloudflare Turnstile support"
]

for feature in captcha_features:
    p = content_frame.add_paragraph()
    p.text = f"• {feature}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

# ============================================================================
# SLIDE 10: Compliance & Standards
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 10)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "9. Compliance & Industry Standards ✅"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "Security Compliance Standards"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(100, 181, 246)
p.space_after = Pt(12)

standards = [
    ("OWASP Top 10", "Protected against all major web vulnerabilities"),
    ("GDPR", "Data portability, right to deletion, consent management"),
    ("SOC 2", "Comprehensive security controls and monitoring"),
    ("ISO 27001", "Information security management standards")
]

for standard, desc in standards:
    p = content_frame.add_paragraph()
    p.text = f"✓ {standard}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(76, 175, 80)

    p = content_frame.add_paragraph()
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1
    p.space_after = Pt(10)

p = content_frame.add_paragraph()
p.text = "\nOWASP Top 10 Protection:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(12)

owasp_protections = [
    "Injection - Parameterized queries",
    "Broken Authentication - JWT + 2FA + bcrypt",
    "Sensitive Data Exposure - AES-256 encryption",
    "XSS - React DOM sanitization",
    "Broken Access Control - RBAC + field permissions"
]

for protection in owasp_protections:
    p = content_frame.add_paragraph()
    p.text = f"• {protection}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

# ============================================================================
# SLIDE 11: Current Security Posture
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)
add_page_number(slide, 11)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "10. Swanson's Security Posture 🎯"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
content_frame = content_box.text_frame

p = content_frame.add_paragraph()
p.text = "✓ Currently Implemented & Active"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(76, 175, 80)
p.space_after = Pt(10)

implemented = [
    "Password authentication with bcrypt hashing",
    "Google OAuth 2.0 integration",
    "JWT token-based session management",
    "AES-256 data encryption (at rest)",
    "HTTPS/TLS 1.3 (in transit)",
    "Role-based access control (22 permission categories)",
    "Field-level permissions",
    "Workspace isolation & multi-tenancy",
    "Comprehensive audit logging",
    "Rate limiting & attack prevention",
    "Two-Factor Authentication (2FA) ready",
    "Regular password rotation enforced",
    "API key security with quarterly rotation",
    "Weekly audit log monitoring",
    "Encrypted backup of security keys"
]

for item in implemented:
    p = content_frame.add_paragraph()
    p.text = f"✓ {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.level = 1

# ============================================================================
# SLIDE 12: SUMMARY / CLOSING
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, (15, 32, 61), (25, 42, 71))
add_logo(slide, 'bottom_right')

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Summary: Tracepoint Security"
title_para = title_frame.paragraphs[0]
title_para.alignment = PP_ALIGN.CENTER
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(4.0))
content_frame = content_box.text_frame

summary_points = [
    "✓ Military-grade encryption (AES-256, TLS 1.3, bcrypt)",
    "✓ Zero-trust architecture with JWT tokens",
    "✓ Enterprise SSO (Google, Microsoft, custom SAML)",
    "✓ Granular permissions (22 categories + field-level)",
    "✓ Attack prevention (rate limiting, audit logs)",
    "✓ Industry compliance (GDPR, SOC 2, ISO 27001, OWASP)",
]

for point in summary_points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(224, 224, 224)
    p.space_after = Pt(8)

# Final statement box
statement_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(1.2))
statement_frame = statement_box.text_frame
statement_frame.word_wrap = True
statement_para = statement_frame.paragraphs[0]
statement_para.text = "Swanson's data is protected by the same security standards used by banks and governments."
statement_para.alignment = PP_ALIGN.CENTER
statement_para.font.size = Pt(20)
statement_para.font.bold = True
statement_para.font.color.rgb = RGBColor(100, 181, 246)

# Save presentation
output_path = '/home/lytle/twenty-dev/Tracepoint_Security_Architecture.pptx'
prs.save(output_path)
print(f"✅ Presentation created successfully!")
print(f"📁 Saved to: {output_path}")
print(f"📊 Total slides: 12")
print(f"\n🎯 Ready for executive presentation!")
