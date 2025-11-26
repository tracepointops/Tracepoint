#!/usr/bin/env python3
"""
Tracepoint Platform - CAPEX Proposal Presentation
Author: Wayne Lytle, Swanson Industries
Purpose: Executive CAPEX Proposal for Unified Business Operations Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

LOGO_PATH = '/home/lytle/twenty-dev/packages/twenty-front/public/images/logos/swanson logo full (1) (1).png'

def add_gradient_background(slide, start_color=(10, 25, 47), end_color=(30, 60, 114)):
    """Enhanced security-themed gradient background"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    fill.gradient_stops[0].color.rgb = RGBColor(*start_color)
    fill.gradient_stops[1].color.rgb = RGBColor(*end_color)

def add_logo(slide, position='top_left'):
    """Add Swanson logo"""
    if not os.path.exists(LOGO_PATH):
        return
    if position == 'top_left':
        left, top, height = Inches(0.5), Inches(0.3), Inches(0.8)
    else:  # bottom_right
        left, top, height = Inches(8.0), Inches(6.8), Inches(0.6)
    try:
        slide.shapes.add_picture(LOGO_PATH, left, top, height=height)
    except:
        pass

print("🚀 Building Tracepoint CAPEX Presentation - Part 1 (Slides 1-5)")

# ============================================================================
# SLIDE 1: COVER
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, (5, 15, 35), (20, 40, 85))
add_logo(slide, 'top_left')

title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
tf = title_box.text_frame
tf.text = "Tracepoint"
p = tf.paragraphs[0]
p.alignment, p.font.size, p.font.bold, p.font.color.rgb = PP_ALIGN.CENTER, Pt(64), True, RGBColor(255, 255, 255)

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(8), Inches(1.0))
tf = subtitle_box.text_frame
tf.text = "AI-Orchestrated Business Operations Platform"
p = tf.paragraphs[0]
p.alignment, p.font.size, p.font.color.rgb = PP_ALIGN.CENTER, Pt(28), RGBColor(100, 181, 246)

tagline_box = slide.shapes.add_textbox(Inches(2), Inches(5.2), Inches(6), Inches(0.6))
tf = tagline_box.text_frame
tf.text = "From Opportunity to Invoice - One Unified Workspace"
p = tf.paragraphs[0]
p.alignment, p.font.size, p.font.color.rgb = PP_ALIGN.CENTER, Pt(18), RGBColor(200, 200, 200)

author_box = slide.shapes.add_textbox(Inches(2), Inches(6.3), Inches(6), Inches(0.8))
tf = author_box.text_frame
tf.text = "CAPEX Proposal\nPresented by: Wayne Lytle\nSwanson Industries"
p = tf.paragraphs[0]
p.alignment, p.font.size, p.font.color.rgb = PP_ALIGN.CENTER, Pt(14), RGBColor(169, 169, 169)

print("✓ Slide 1: Cover")

# ============================================================================
# SLIDE 2: EXECUTIVE OVERVIEW
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "What Tracepoint Is"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame
cf.word_wrap = True

# Main description
p = cf.add_paragraph()
p.text = "Unified Business Operations Platform"
p.font.size, p.font.bold, p.font.color.rgb = Pt(24), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Tracepoint merges enterprise CRM with operational depth into a single, AI-orchestrated workspace that centralizes:"
p.font.size, p.font.color.rgb = Pt(16), RGBColor(255, 255, 255)
p.space_after = Pt(10)

modules = [
    "Quoting & Pricing",
    "Orders & Purchase Orders",
    "Work Orders & Project Management",
    "Shipments & Logistics",
    "Invoicing & Financial Intelligence",
    "CRM Pipelines & Customer Management"
]

for module in modules:
    p = cf.add_paragraph()
    p.text = f"✓ {module}"
    p.font.size, p.font.color.rgb, p.level = Pt(15), RGBColor(76, 175, 80), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Powered by SIA - Your Adaptive AI Assistant"
p.font.size, p.font.bold, p.font.color.rgb = Pt(18), True, RGBColor(255, 193, 7)
p.space_before = Pt(15)

p = cf.add_paragraph()
p.text = "All data flows from QAD ERP through automated pipelines, ensuring real-time accuracy while Tracepoint becomes your daily operational workspace."
p.font.size, p.font.color.rgb = Pt(14), RGBColor(224, 224, 224)
p.level = 1

print("✓ Slide 2: Executive Overview")

# ============================================================================
# SLIDE 3: ARCHITECTURE & DATA FLOW
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Architecture & Data Flow"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Dual Interface Ecosystem"
p.font.size, p.font.bold, p.font.color.rgb = Pt(22), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

architecture = [
    ("QAD ERP", "Authoritative source of truth for enterprise data"),
    ("Tracepoint", "Daily operational workspace for all contributors"),
    ("Automated Pipeline", "CSV reports flow to PostgreSQL + Redis for real-time access")
]

for component, desc in architecture:
    p = cf.add_paragraph()
    p.text = component
    p.font.size, p.font.bold, p.font.color.rgb = Pt(18), True, RGBColor(255, 255, 255)
    p.space_before = Pt(8)

    p = cf.add_paragraph()
    p.text = f"→ {desc}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(230, 230, 230), 1
    p.space_after = Pt(6)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Data Backbone Benefits:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(18), True, RGBColor(100, 181, 246)
p.space_before = Pt(15)

benefits = [
    "Every module pulls from the same normalized data layer",
    "Ensures consistency and complete traceability",
    "Real-time visibility across all operations",
    "Single source of truth with zero data duplication"
]

for benefit in benefits:
    p = cf.add_paragraph()
    p.text = f"• {benefit}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1

print("✓ Slide 3: Architecture")

# ============================================================================
# SLIDE 4: AI ORCHESTRATION TEAM
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "AI Orchestration Team"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Meet SIA & The Intelligent Agent Team"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Users interact with SIA - a professional, adaptive assistant. Behind the scenes, specialized agents orchestrate every workflow:"
p.font.size, p.font.color.rgb = Pt(14), RGBColor(224, 224, 224)
p.space_after = Pt(12)

agents = [
    ("SIA", "Smart Intelligent Assistant", "Your voice, guide, and interface to the AI team"),
    ("ADA", "Language ↔ SQL Translator", "Converts questions into queries and data into insights"),
    ("DEX", "Logic ↔ App ↔ Render", "Transforms UI interactions into database operations"),
    ("CID", "Context & Documentation", "Answers questions from company knowledge and ERP docs"),
    ("ARC", "Visualization AI", "Renders charts, dashboards, scorecards automatically"),
    ("DEV", "Admin AI", "Manages users, permissions, and company profiles")
]

for agent_name, agent_title, agent_role in agents:
    p = cf.add_paragraph()
    p.text = f"{agent_name} - {agent_title}"
    p.font.size, p.font.bold, p.font.color.rgb = Pt(15), True, RGBColor(255, 193, 7)

    p = cf.add_paragraph()
    p.text = agent_role
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(230, 230, 230), 1
    p.space_after = Pt(6)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Result: Contributors work faster with less training, while AI ensures accuracy and compliance."
p.font.size, p.font.italic, p.font.color.rgb = Pt(14), True, RGBColor(76, 175, 80)
p.space_before = Pt(10)

print("✓ Slide 4: AI Team")

# ============================================================================
# SLIDE 5: QUOTE MANAGEMENT MODULE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Quote Management"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Enterprise Quote-to-Cash Workflow Automation"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Streamline the entire quoting lifecycle with intelligent automation and real-time visibility."
p.font.size, p.font.color.rgb = Pt(15), RGBColor(255, 255, 255)
p.space_after = Pt(12)

features = [
    "Multi-stage approval workflows with configurable business rules",
    "Automated pricing calculations and version control",
    "Track status from request → approval → acceptance → order conversion",
    "Real-time analytics: win rates, quote values, bottleneck identification",
    "Integration with financial systems for accurate pricing",
    "Complete audit trails for regulatory compliance"
]

for feature in features:
    p = cf.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Key Benefits:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 193, 7)
p.space_before = Pt(12)

benefits = [
    "Faster quote turnaround with automated approvals",
    "Higher win rates through predictive analytics",
    "Reduced errors with template-based generation",
    "Complete visibility for sales leadership"
]

for benefit in benefits:
    p = cf.add_paragraph()
    p.text = f"→ {benefit}"
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(76, 175, 80), 1

print("✓ Slide 5: Quote Management")

# Save Part 1
output_path = '/home/lytle/twenty-dev/Tracepoint_CAPEX_Part1.pptx'
prs.save(output_path)
print(f"\n✅ Part 1 Complete (Slides 1-5)")
print(f"📁 Saved: {output_path}")
print(f"\n🔄 Ready for Part 2 (Slides 6-10)")

# ============================================================================
# PART 2: SLIDES 6-10
# ============================================================================
print("\n🚀 Building Part 2 (Slides 6-10)")

# ============================================================================
# SLIDE 6: SHIPMENT MANAGEMENT
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Shipment Management"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "End-to-End Logistics Command Center"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Transform logistics operations with real-time visibility across the entire supply chain."
p.font.size, p.font.color.rgb = Pt(15), RGBColor(255, 255, 255)
p.space_after = Pt(12)

features = [
    "Live GPS tracking with geofencing alerts and route optimization",
    "Automated driver dispatch with intelligent load balancing",
    "Monitor from pickup → delivery with milestone tracking",
    "Predictive delivery windows and exception alerts",
    "Mobile-first design for drivers with digital proof of delivery",
    "Complete audit logging for regulatory compliance"
]

for feature in features:
    p = cf.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Performance Metrics:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 193, 7)
p.space_before = Pt(12)

metrics = [
    "Driver efficiency tracking and performance analytics",
    "SLA monitoring and customer satisfaction scoring",
    "Chain-of-custody documentation for regulated shipments",
    "Integration with third-party carriers and freight systems"
]

for metric in metrics:
    p = cf.add_paragraph()
    p.text = f"→ {metric}"
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(76, 175, 80), 1

print("✓ Slide 6: Shipment Management")

# ============================================================================
# SLIDE 7: PROJECT MANAGEMENT
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Project Management"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Collaborative Work Orchestration Platform"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Drive project success with enterprise-grade task management and team collaboration."
p.font.size, p.font.color.rgb = Pt(15), RGBColor(255, 255, 255)
p.space_after = Pt(12)

features = [
    "Unlimited task nesting with dependencies and milestone tracking",
    "Intelligent workload distribution across teams",
    "Automated workflows for notifications, approvals, escalations",
    "Real-time dashboards showing project health and timeline adherence",
    "Gantt charts and critical path analysis",
    "Time tracking with billable hours management"
]

for feature in features:
    p = cf.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Collaboration Features:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 193, 7)
p.space_before = Pt(12)

collab = [
    "Threaded comments and @mentions for seamless communication",
    "File attachments and document sharing",
    "Resource capacity planning and cross-project portfolio management",
    "Integration with calendars and communication platforms"
]

for item in collab:
    p = cf.add_paragraph()
    p.text = f"→ {item}"
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(76, 175, 80), 1

print("✓ Slide 7: Project Management")

# ============================================================================
# SLIDE 8: FINANCIAL INTELLIGENCE (SIA)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Financial Intelligence (SIA)"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "AI-Powered Business Intelligence Assistant"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Ask complex business questions in natural language - SIA delivers instant, data-backed insights."
p.font.size, p.font.color.rgb = Pt(15), RGBColor(255, 255, 255)
p.space_after = Pt(12)

capabilities = [
    "Natural language queries against live financial data",
    "Automatically generates and executes SQL queries",
    "Executive summaries with actionable recommendations",
    "Learns your terminology, metrics, and decision patterns",
    "Proactive alerting for trends and anomalies",
    "Customer segmentation and lifetime value calculations"
]

for cap in capabilities:
    p = cf.add_paragraph()
    p.text = f"✓ {cap}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Business Impact:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 193, 7)
p.space_before = Pt(12)

impact = [
    "Democratize data access - no SQL knowledge required",
    "Revenue forecasting and predictive analytics",
    "Competitive benchmarking and market intelligence",
    "Secure, role-based access with row-level security"
]

for item in impact:
    p = cf.add_paragraph()
    p.text = f"→ {item}"
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(76, 175, 80), 1

print("✓ Slide 8: Financial Intelligence")

# ============================================================================
# SLIDE 9: DOCUMENT PROCESSING
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Document Processing"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Intelligent Document Automation Engine"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Eliminate manual paperwork with AI-powered document generation and processing."
p.font.size, p.font.color.rgb = Pt(15), RGBColor(255, 255, 255)
p.space_after = Pt(12)

features = [
    "Auto-generate bills of lading, delivery tickets, packing slips",
    "Template-based generation ensures brand consistency",
    "Digital signature capture with tamper-proof timestamps",
    "OCR and AI extract data from incoming documents",
    "Secure cloud storage with version control and audit trails",
    "Automated document routing and approval workflows"
]

for feature in features:
    p = cf.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Compliance & Security:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 193, 7)
p.space_before = Pt(12)

compliance = [
    "Legal proof of delivery and acceptance",
    "Encryption at rest and in transit",
    "Compliance-ready retention policies and legal holds",
    "Complete document history for audit investigations"
]

for item in compliance:
    p = cf.add_paragraph()
    p.text = f"→ {item}"
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(76, 175, 80), 1

print("✓ Slide 9: Document Processing")

# ============================================================================
# SLIDE 10: USER & ROLE MANAGEMENT
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "User & Role Management"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Zero-Trust Security and Access Control"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Enterprise-grade security with granular role-based access control (RBAC)."
p.font.size, p.font.color.rgb = Pt(15), RGBColor(255, 255, 255)
p.space_after = Pt(12)

features = [
    "Custom roles with precise permissions at object, field, and record levels",
    "Department-based data segregation",
    "Automated provisioning/de-provisioning with HR system integration",
    "Multi-factor authentication and session management",
    "Single Sign-On (SSO) with enterprise identity providers",
    "Real-time security monitoring and threat detection"
]

for feature in features:
    p = cf.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Compliance & Audit:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 193, 7)
p.space_before = Pt(12)

audit = [
    "Comprehensive audit logs for every login and data access",
    "Privileged access management for admin functions",
    "IP whitelisting and geographic restrictions",
    "Security investigations and compliance reporting"
]

for item in audit:
    p = cf.add_paragraph()
    p.text = f"→ {item}"
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(76, 175, 80), 1

print("✓ Slide 10: User & Role Management")

# Save Part 2
prs.save(output_path)
print(f"\n✅ Part 2 Complete (Slides 6-10)")
print(f"📁 Updated: {output_path}")
print(f"\n🔄 Ready for Part 3 (Slides 11-15)")

# ============================================================================
# PART 3: SLIDES 11-15 (FINAL)
# ============================================================================
print("\n🚀 Building Part 3 (Slides 11-15 - FINAL)")

# ============================================================================
# SLIDE 11: CALENDAR INTEGRATION
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Calendar Integration"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Unified Scheduling and Resource Coordination"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Bi-directional Google Calendar integration synchronizes operations across the enterprise."
p.font.size, p.font.color.rgb = Pt(15), RGBColor(255, 255, 255)
p.space_after = Pt(12)

features = [
    "Automated scheduling of field service appointments and delivery windows",
    "Intelligent conflict detection and resolution",
    "Coordinate driver schedules and equipment availability in real-time",
    "Mobile notifications for schedule changes and updates",
    "Customer self-service scheduling portals",
    "Integration with project management for task deadline alignment"
]

for feature in features:
    p = cf.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Performance Optimization:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 193, 7)
p.space_before = Pt(12)

benefits = [
    "Analytics identify scheduling inefficiencies",
    "Optimize resource utilization and on-time performance",
    "Improve customer satisfaction with reliable appointment windows"
]

for benefit in benefits:
    p = cf.add_paragraph()
    p.text = f"→ {benefit}"
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(76, 175, 80), 1

print("✓ Slide 11: Calendar Integration")

# ============================================================================
# SLIDE 12: ANALYTICS & REPORTING
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Analytics & Reporting"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Real-Time Business Intelligence Platform"
p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, RGBColor(100, 181, 246)
p.space_after = Pt(12)

p = cf.add_paragraph()
p.text = "Transform data into actionable insights with enterprise-grade analytics."
p.font.size, p.font.color.rgb = Pt(15), RGBColor(255, 255, 255)
p.space_after = Pt(12)

capabilities = [
    "Pre-built dashboards for instant KPI visibility across all operations",
    "Self-service report builder with drag-and-drop interface",
    "AI-driven insights identify trends and optimization opportunities",
    "Scheduled reports delivered via email or messaging platforms",
    "Data visualization transforms datasets into intuitive charts and heatmaps",
    "Integration with Looker Studio, Power BI, and Tableau"
]

for cap in capabilities:
    p = cf.add_paragraph()
    p.text = f"✓ {cap}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(224, 224, 224), 1
    p.space_after = Pt(4)

cf.add_paragraph()
p = cf.add_paragraph()
p.text = "Advanced Features:"
p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 193, 7)
p.space_before = Pt(12)

advanced = [
    "Predictive analytics and machine learning models",
    "Automated anomaly detection and alerting",
    "Data governance and certified metrics repository"
]

for item in advanced:
    p = cf.add_paragraph()
    p.text = f"→ {item}"
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(76, 175, 80), 1

print("✓ Slide 12: Analytics & Reporting")

# ============================================================================
# SLIDE 13: ENTERPRISE PLATFORM ADVANTAGES
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Enterprise Platform Advantages"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

advantages = [
    ("Scalability", "Cloud-native architecture supports unlimited users and data growth without performance degradation"),
    ("Security", "Enterprise-grade encryption, RBAC, and comprehensive audit logging meet SOC 2 and ISO 27001 requirements"),
    ("Integration", "RESTful APIs, webhooks, and pre-built connectors enable seamless ERP and legacy system integration"),
    ("Customization", "Configurable workflows and custom fields adapt to unique business processes without custom development"),
    ("Reliability", "99.9% uptime SLA, automated backups, and multi-region redundancy ensure business continuity"),
    ("Support", "24/7 technical support and comprehensive training programs accelerate adoption and maximize ROI")
]

for title, desc in advantages:
    p = cf.add_paragraph()
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(18), True, RGBColor(255, 193, 7)
    p.space_before = Pt(6)
    
    p = cf.add_paragraph()
    p.text = desc
    p.font.size, p.font.color.rgb, p.level = Pt(13), RGBColor(230, 230, 230), 1
    p.space_after = Pt(8)

print("✓ Slide 13: Enterprise Advantages")

# ============================================================================
# SLIDE 14: STRATEGIC BUSINESS VALUE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
tf = title_box.text_frame
tf.text = "Strategic Business Value"
tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = Pt(40), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.8))
cf = content_box.text_frame

p = cf.add_paragraph()
p.text = "Why Swanson Industries Needs Tracepoint"
p.font.size, p.font.bold, p.font.color.rgb = Pt(22), True, RGBColor(100, 181, 246)
p.space_after = Pt(15)

value_props = [
    ("Accelerated Decision-Making", "AI-driven insights enable faster, better business decisions"),
    ("Enhanced Operational Efficiency", "Workflow automation significantly reduces manual efforts"),
    ("Improved Transparency", "Real-time visibility drives clear ownership and accountability"),
    ("Contributor Empowerment", "Adaptive AI guidance minimizes training and maximizes productivity"),
    ("Future-Proof Growth", "Extensible platform scales with business needs"),
    ("Competitive Advantage", "Unified operations platform positions Swanson ahead of competitors")
]

for title, value in value_props:
    p = cf.add_paragraph()
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(16), True, RGBColor(255, 255, 255)
    
    p = cf.add_paragraph()
    p.text = f"→ {value}"
    p.font.size, p.font.color.rgb, p.level = Pt(14), RGBColor(76, 175, 80), 1
    p.space_after = Pt(8)

print("✓ Slide 14: Strategic Value")

# ============================================================================
# SLIDE 15: SUMMARY (WITH LOGO BOTTOM RIGHT)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, (5, 15, 35), (20, 40, 85))
add_logo(slide, 'bottom_right')

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.9))
tf = title_box.text_frame
tf.text = "Tracepoint: Your Unified Operations Platform"
tf.paragraphs[0].alignment, tf.paragraphs[0].font.size, tf.paragraphs[0].font.bold, tf.paragraphs[0].font.color.rgb = PP_ALIGN.CENTER, Pt(42), True, RGBColor(255, 255, 255)

content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(7.6), Inches(4.2))
cf = content_box.text_frame

summary_points = [
    "AI-orchestrated platform merging CRM with operational modules",
    "Single workspace from opportunity to invoice",
    "SIA and intelligent agent team automate workflows",
    "Real-time data flow from QAD ERP to daily operations",
    "Enterprise-grade security, scalability, and compliance",
    "Proven ROI through efficiency, transparency, and growth"
]

for point in summary_points:
    p = cf.add_paragraph()
    p.text = f"✓ {point}"
    p.font.size, p.font.color.rgb = Pt(17), RGBColor(224, 224, 224)
    p.space_after = Pt(10)

# Final statement box
statement_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(7), Inches(0.9))
tf = statement_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Tracepoint delivers enterprise-grade operational excellence through intelligent automation, real-time visibility, and AI-powered decision support."
p.alignment, p.font.size, p.font.bold, p.font.color.rgb = PP_ALIGN.CENTER, Pt(16), True, RGBColor(100, 181, 246)

print("✓ Slide 15: Summary")

# Save final presentation
prs.save(output_path)
print(f"\n" + "="*60)
print(f"✅ COMPLETE! All 15 slides generated successfully!")
print(f"📁 Final presentation: {output_path}")
print(f"📊 Total slides: 15")
print(f"💼 Ready for executive CAPEX proposal presentation")
print(f"="*60)
